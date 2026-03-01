"""Parse uploaded documents into plain text."""

from __future__ import annotations

import hashlib
import io
import logging
from typing import Callable

TYPE_PROGRESS_CB = Callable[[int, int], None]  # (current, total); total=0 means unknown
import re
import shutil
import subprocess
import tempfile
import zipfile
from collections import defaultdict
from pathlib import Path

import pdfplumber
from docx import Document as DocxDocument
from PIL import Image, UnidentifiedImageError
from pptx import Presentation
import trafilatura

from app.config import ALLOW_REMOTE_URL_INGEST, LOCAL_ONLY, OUTPUTS_DIR, get_ocr_settings, get_vision_ingest_settings

logger = logging.getLogger(__name__)

OCR_DEFAULT_ENABLED = True
OCR_DEFAULT_MODE = "fast"
OCR_DEFAULT_LANG = "rus+eng"
OCR_DEFAULT_MIN_CHARS = 8
OCR_DEFAULT_MAX_PDF_PAGES = 40
OCR_DEFAULT_MAX_DOCX_IMAGES = 40
_TESSERACT_OK: bool | None = None
PDF_MARGIN_SCAN_LINES = 3
PDF_MARGIN_REPEAT_MIN_PAGES = 3
PDF_MARGIN_CANDIDATE_MAX_LEN = 140
PDF_TABLE_MAX_PER_PAGE = 6
PDF_TABLE_MAX_ROWS = 24
PDF_TABLE_MAX_COLS = 10
PDF_TABLE_MIN_FILLED_CELLS = 4
PDF_FIGURE_CAPTION_MAX_PER_PAGE = 8
PPTX_TABLE_MAX_PER_SLIDE = 8
PPTX_TABLE_MAX_ROWS = 24
PPTX_TABLE_MAX_COLS = 10
PPTX_FIGURE_CAPTION_MAX_PER_SLIDE = 8
# Vision ingest: skip very small images (icons, thumbnails)
VISION_MIN_SIDE_PX = 64
VISION_MIN_PIXELS = 2500
# PDF: skip images in header/footer (fraction of page height from top/bottom)
VISION_PDF_HEADER_FRAC = 0.12
VISION_PDF_FOOTER_FRAC = 0.12


def _vision_skip_small_image(width: int, height: int) -> bool:
    """Return True if image is too small to send to vision LLM."""
    if width <= 0 or height <= 0:
        return True
    if width < VISION_MIN_SIDE_PX and height < VISION_MIN_SIDE_PX:
        return True
    if width * height < VISION_MIN_PIXELS:
        return True
    return False


def _vision_pdf_skip_header_footer(rect_y0: float, rect_y1: float, page_height: float) -> bool:
    """Return True if image rect is entirely in header or footer zone (skip it)."""
    if page_height <= 0:
        return False
    top_cut = page_height * VISION_PDF_HEADER_FRAC
    bottom_cut = page_height * (1 - VISION_PDF_FOOTER_FRAC)
    # In PDF coordinates y often grows upward; in PyMuPDF Rect, y0 is top, y1 is bottom in page coords
    if rect_y1 <= top_cut:
        return True  # entirely in header
    if rect_y0 >= bottom_cut:
        return True  # entirely in footer
    return False


_PDF_MARKER_RE = re.compile(r"^\[(?:OCR\s+)?PDF page\s+\d+\]$", flags=re.IGNORECASE)
_DOCX_HEADING_STYLE_RE = re.compile(r"(?:^|\s)(?:heading|заголовок)\s*([1-6])?\b", flags=re.IGNORECASE)
_TABLE_CAPTION_RE = re.compile(
    r"^\s*(?:табл(?:\.|ица)?|table)\s*(?:\d+[.\d]*|[IVXLCM]+)?\s*[:.)-]?\s+.+$",
    flags=re.IGNORECASE,
)
_FIGURE_CAPTION_RE = re.compile(
    r"^\s*(?:рис(?:\.|унок)?|fig(?:\.|ure)?)\s*(?:\d+[.\d]*|[IVXLCM]+)?\s*[:.)-]?\s+.+$",
    flags=re.IGNORECASE,
)
_CAPTION_LINE_RE = re.compile(r"^\s*Подпись:\s*(.+?)\s*$", flags=re.IGNORECASE)
_ANCHOR_LINE_RE = re.compile(r"^\s*Якорь:\s*(.+?)\s*$", flags=re.IGNORECASE)
_MD_HEADING_LINE_RE = re.compile(r"^\s{0,3}#{1,6}\s+(.+?)\s*$")


OFFICE_WORD_SUFFIXES = {".docx", ".doc", ".rtf", ".odt", ".otd"}
OFFICE_PRESENTATION_SUFFIXES = {".pptx", ".ppt"}
DJVU_SUFFIXES = {".djvu", ".djv", ".djvy"}
PREVIEWABLE_CONVERT_SUFFIXES = OFFICE_WORD_SUFFIXES | OFFICE_PRESENTATION_SUFFIXES | DJVU_SUFFIXES


def parse_file(path: Path, *, progress_cb: TYPE_PROGRESS_CB | None = None) -> str:
    """Return plain-text content from a PDF, office document, DJVU, text, or saved HTML.
    If progress_cb is set, it is called during vision image processing as progress_cb(current, total); total=0 means unknown.
    """
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _parse_pdf(path, progress_cb=progress_cb)
    if suffix == ".docx":
        return _parse_docx(path, progress_cb=progress_cb)
    if suffix == ".pptx":
        return _parse_pptx(path, progress_cb=progress_cb)
    if suffix in {".doc", ".rtf", ".odt", ".otd", ".ppt"}:
        return _parse_via_preview_pdf(path, progress_cb=progress_cb)
    if suffix in DJVU_SUFFIXES:
        return _parse_via_preview_pdf(path, progress_cb=progress_cb)
    if suffix in (".txt", ".md", ".rst"):
        return path.read_text(encoding="utf-8", errors="replace")
    if suffix in (".html", ".htm"):
        return _parse_html(path.read_text(encoding="utf-8", errors="replace"))
    raise ValueError(f"Unsupported file type: {suffix}")


def parse_url(url: str) -> str:
    """Fetch and extract text from a URL using trafilatura."""
    if LOCAL_ONLY and not ALLOW_REMOTE_URL_INGEST:
        raise ValueError("LOCAL_ONLY режим: импорт по URL отключен. Загрузите файл напрямую.")
    downloaded = trafilatura.fetch_url(url)
    if downloaded is None:
        raise ValueError(f"Could not fetch URL: {url}")
    text = trafilatura.extract(downloaded)
    if not text:
        raise ValueError(f"Could not extract text from URL: {url}")
    return text


def ensure_preview_pdf(path: Path, *, document_id: str | None = None) -> Path:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return path
    if suffix in OFFICE_WORD_SUFFIXES | OFFICE_PRESENTATION_SUFFIXES:
        return _convert_office_to_pdf(path, document_id=document_id)
    if suffix in DJVU_SUFFIXES:
        return _convert_djvu_to_pdf(path, document_id=document_id)
    raise ValueError(f"Preview is not supported for {suffix}")


def _parse_via_preview_pdf(path: Path, *, progress_cb: TYPE_PROGRESS_CB | None = None) -> str:
    with tempfile.TemporaryDirectory(prefix="ingest-preview-") as tmp:
        pdf_path = ensure_preview_pdf(path, document_id=Path(tmp, path.stem).name)
        if pdf_path.parent == OUTPUTS_DIR:
            temp_pdf = Path(tmp) / f"{path.stem}.pdf"
            shutil.copy2(pdf_path, temp_pdf)
            pdf_path = temp_pdf
        return _parse_pdf(pdf_path, progress_cb=progress_cb)


def _preview_output_path(document_id: str | None, source_path: Path) -> Path:
    stem = str(document_id or source_path.stem or "preview").strip() or "preview"
    return OUTPUTS_DIR / f"{stem}_preview.pdf"


def _is_fresh_artifact(source_path: Path, artifact_path: Path) -> bool:
    if not artifact_path.exists() or artifact_path.stat().st_size <= 0:
        return False
    try:
        return artifact_path.stat().st_mtime_ns >= source_path.stat().st_mtime_ns
    except OSError:
        return False


def _resolve_soffice_binary() -> str:
    for candidate in ("soffice", "libreoffice"):
        found = shutil.which(candidate)
        if found:
            return found
    raise ValueError("LibreOffice headless не установлен в backend-образе")


def _run_subprocess(args: list[str], *, timeout: int = 180, error_prefix: str = "Команда завершилась с ошибкой") -> subprocess.CompletedProcess:
    try:
        return subprocess.run(args, check=True, capture_output=True, text=True, timeout=timeout)
    except FileNotFoundError as exc:
        raise ValueError(f"{error_prefix}: бинарник не найден ({args[0]})") from exc
    except subprocess.TimeoutExpired as exc:
        raise ValueError(f"{error_prefix}: превышен таймаут {timeout}s") from exc
    except subprocess.CalledProcessError as exc:
        stderr = (exc.stderr or exc.stdout or "").strip()
        detail = f": {stderr}" if stderr else ""
        raise ValueError(f"{error_prefix}{detail}") from exc


def _convert_office_to_pdf(path: Path, *, document_id: str | None = None) -> Path:
    out_path = _preview_output_path(document_id, path)
    if document_id and _is_fresh_artifact(path, out_path):
        return out_path
    soffice = _resolve_soffice_binary()
    with tempfile.TemporaryDirectory(prefix="soffice-preview-") as tmp:
        tmp_dir = Path(tmp)
        work_in = tmp_dir / "in"
        work_out = tmp_dir / "out"
        work_in.mkdir(parents=True, exist_ok=True)
        work_out.mkdir(parents=True, exist_ok=True)
        src_copy = work_in / path.name
        shutil.copy2(path, src_copy)
        _run_subprocess([
            soffice,
            "--headless",
            "--nologo",
            "--nodefault",
            "--nofirststartwizard",
            "--nolockcheck",
            "--convert-to",
            "pdf",
            "--outdir",
            str(work_out),
            str(src_copy),
        ], error_prefix=f"Не удалось конвертировать {path.suffix} в PDF")
        candidates = sorted(work_out.glob("*.pdf"))
        if not candidates:
            raise ValueError(f"LibreOffice не создал preview PDF для {path.name}")
        pdf_path = candidates[0]
        if document_id:
            OUTPUTS_DIR.mkdir(parents=True, exist_ok=True)
            tmp_out = out_path.with_suffix(out_path.suffix + ".tmp")
            shutil.copy2(pdf_path, tmp_out)
            tmp_out.replace(out_path)
            return out_path
        return pdf_path


def _convert_djvu_to_pdf(path: Path, *, document_id: str | None = None) -> Path:
    out_path = _preview_output_path(document_id, path)
    if document_id and _is_fresh_artifact(path, out_path):
        return out_path
    target = out_path if document_id else Path(tempfile.mkdtemp(prefix="djvu-preview-")) / f"{path.stem}.pdf"
    target.parent.mkdir(parents=True, exist_ok=True)
    _run_subprocess(["ddjvu", "-format=pdf", str(path), str(target)], error_prefix=f"Не удалось конвертировать {path.suffix} в PDF")
    if not target.exists() or target.stat().st_size <= 0:
        raise ValueError(f"ddjvu не создал preview PDF для {path.name}")
    return target


# -- private helpers -------------------------------------------------------

def _ocr_runtime_settings() -> dict:
    try:
        cfg = get_ocr_settings()
    except Exception:
        cfg = {}
    mode = str(cfg.get("mode", OCR_DEFAULT_MODE)).strip().lower()
    if mode not in {"fast", "accurate"}:
        mode = OCR_DEFAULT_MODE
    def _safe_int(value, default):
        try:
            return int(value)
        except (TypeError, ValueError):
            return int(default)

    return {
        "enabled": bool(cfg.get("enabled", OCR_DEFAULT_ENABLED)),
        "mode": mode,
        "lang": str(cfg.get("lang", OCR_DEFAULT_LANG)).strip() or OCR_DEFAULT_LANG,
        "min_chars": max(1, _safe_int(cfg.get("min_chars", OCR_DEFAULT_MIN_CHARS), OCR_DEFAULT_MIN_CHARS)),
        "max_pdf_pages": max(1, _safe_int(cfg.get("max_pdf_pages", OCR_DEFAULT_MAX_PDF_PAGES), OCR_DEFAULT_MAX_PDF_PAGES)),
        "max_docx_images": max(1, _safe_int(cfg.get("max_docx_images", OCR_DEFAULT_MAX_DOCX_IMAGES), OCR_DEFAULT_MAX_DOCX_IMAGES)),
    }

def _count_pdf_vision_candidates(path: Path, max_total: int) -> int:
    """Count how many PDF images will be sent to vision (for progress N/M). PyMuPDF only."""
    try:
        import fitz
    except ImportError:
        return 0
    try:
        doc = fitz.open(str(path))
        seen_hashes: set[str] = set()
        count = 0
        for page_no in range(len(doc)):
            if count >= max_total:
                break
            page = doc[page_no]
            page_height = page.rect.height
            for img in doc.get_page_images(page_no):
                if count >= max_total:
                    break
                xref = img[0]
                try:
                    rects = page.get_image_rects(xref)
                    if rects and _vision_pdf_skip_header_footer(rects[0].y0, rects[0].y1, page_height):
                        continue
                except Exception:
                    pass
                info = doc.extract_image(xref)
                if not info:
                    continue
                raw = info.get("image")
                w = info.get("width", 0)
                h = info.get("height", 0)
                if not raw or _vision_skip_small_image(w, h):
                    continue
                hsh = hashlib.sha256(raw).hexdigest()
                if hsh not in seen_hashes:
                    seen_hashes.add(hsh)
                    count += 1
            try:
                blocks = page.get_text("dict", flags=0).get("blocks") or []
            except Exception:
                blocks = []
            for blk in blocks:
                if count >= max_total:
                    break
                if blk.get("type") != 1 or "image" not in blk:
                    continue
                raw_inline = blk.get("image")
                if not raw_inline or len(raw_inline) < 100:
                    continue
                bbox = blk.get("bbox")
                if bbox and len(bbox) >= 4 and page_height > 0:
                    if _vision_pdf_skip_header_footer(bbox[1], bbox[3], page_height):
                        continue
                try:
                    pil_img = Image.open(io.BytesIO(raw_inline))
                    iw, ih = pil_img.size
                    pil_img.close()
                except Exception:
                    iw, ih = 0, 0
                if _vision_skip_small_image(iw, ih):
                    continue
                hsh = hashlib.sha256(raw_inline).hexdigest()
                if hsh not in seen_hashes:
                    seen_hashes.add(hsh)
                    count += 1
        doc.close()
        return count
    except Exception:
        return 0


def _extract_pdf_vision_descriptions(
    path: Path,
    *,
    progress_cb: TYPE_PROGRESS_CB | None = None,
) -> dict[int, list[str]]:
    """Extract embedded images from PDF pages, describe via vision LLM; return page_idx -> list of description blocks.
    Prefers PyMuPDF (exact embedded image bytes, no crop). Falls back to pdfplumber crop if PyMuPDF unavailable.
    """
    vision_cfg = get_vision_ingest_settings()
    logger.info("Vision ingest (PDF): entry for %s (enabled=%s)", path.name, vision_cfg.get("enabled"))
    if not vision_cfg.get("enabled"):
        return {}
    max_total = max(1, int(vision_cfg.get("max_images_per_document", 20)))
    logger.info("Vision ingest (PDF): enabled, max %s images, scanning %s", max_total, path.name)
    page_vision_blocks: dict[int, list[str]] = {}
    try:
        from app.services import llm_service
    except ImportError:
        return {}

    # 1) Prefer PyMuPDF: extract embedded images by xref (no crop, finds all images)
    try:
        import fitz
    except ImportError:
        fitz = None
        logger.debug("Vision ingest (PDF): PyMuPDF not installed, using pdfplumber")
    if fitz is not None:
        try:
            total_to_process = _count_pdf_vision_candidates(path, max_total)
            doc = fitz.open(str(path))
            xref_to_desc: dict[int, str] = {}
            image_hash_to_desc: dict[str, str] = {}  # dedupe: same image as xref and inline only sent once
            total_described = 0
            for page_no in range(len(doc)):
                if total_described >= max_total:
                    break
                page_idx = page_no + 1
                page = doc[page_no]
                page_height = page.rect.height
                imglist = doc.get_page_images(page_no)
                for img_idx, img in enumerate(imglist):
                    if total_described >= max_total:
                        break
                    xref = img[0]
                    try:
                        rects = page.get_image_rects(xref)
                        if rects and _vision_pdf_skip_header_footer(rects[0].y0, rects[0].y1, page_height):
                            logger.debug("PDF vision (PyMuPDF): skip header/footer image page %s xref %s", page_idx, xref)
                            continue
                    except Exception:
                        pass
                    if xref in xref_to_desc:
                        page_vision_blocks.setdefault(page_idx, []).append(
                            f"[PDF page {page_idx}, изображение {img_idx + 1}]\nОписание: {xref_to_desc[xref]}"
                        )
                        continue
                    info = doc.extract_image(xref)
                    if not info:
                        continue
                    raw = info.get("image")
                    w = info.get("width", 0)
                    h = info.get("height", 0)
                    if not raw or _vision_skip_small_image(w, h):
                        if raw:
                            logger.debug("PDF vision (PyMuPDF): skip small image page %s xref %s (%sx%s)", page_idx, xref, w, h)
                        continue
                    content_hash = hashlib.sha256(raw).hexdigest()
                    if content_hash in image_hash_to_desc:
                        desc = image_hash_to_desc[content_hash]
                        xref_to_desc[xref] = desc
                        page_vision_blocks.setdefault(page_idx, []).append(
                            f"[PDF page {page_idx}, изображение {img_idx + 1}]\nОписание: {desc}"
                        )
                        continue
                    ext = (info.get("ext") or "png").lower()
                    mime = "image/jpeg" if ext in ("jpg", "jpeg") else "image/png"
                    desc = llm_service.describe_image_with_vision(raw, mime_type=mime)
                    if desc:
                        xref_to_desc[xref] = desc
                        image_hash_to_desc[content_hash] = desc
                        logger.info("Vision ingest: adding description to PDF page %s image %s (%s chars)", page_idx, img_idx + 1, len(desc))
                        page_vision_blocks.setdefault(page_idx, []).append(
                            f"[PDF page {page_idx}, изображение {img_idx + 1}]\nОписание: {desc}"
                        )
                        total_described += 1
                        if progress_cb:
                            progress_cb(total_described, total_to_process)
                # Inline images (no xref): often diagrams that get_page_images misses; dedupe by content hash
                try:
                    blocks = page.get_text("dict", flags=0).get("blocks") or []
                except Exception:
                    blocks = []
                for blk in blocks:
                    if total_described >= max_total:
                        break
                    if blk.get("type") != 1 or "image" not in blk:
                        continue
                    raw_inline = blk.get("image")
                    if not raw_inline or len(raw_inline) < 100:
                        continue
                    bbox = blk.get("bbox")
                    if bbox and len(bbox) >= 4 and page_height > 0:
                        if _vision_pdf_skip_header_footer(bbox[1], bbox[3], page_height):
                            continue
                    try:
                        pil_img = Image.open(io.BytesIO(raw_inline))
                        iw, ih = pil_img.size
                        pil_img.close()
                    except Exception:
                        iw, ih = 0, 0
                    if _vision_skip_small_image(iw, ih):
                        continue
                    content_hash_inline = hashlib.sha256(raw_inline).hexdigest()
                    if content_hash_inline in image_hash_to_desc:
                        desc_inline = image_hash_to_desc[content_hash_inline]
                        page_vision_blocks.setdefault(page_idx, []).append(
                            f"[PDF page {page_idx}, изображение (схема/диаграмма)]\nОписание: {desc_inline}"
                        )
                        continue
                    mime_inline = "image/png"
                    desc_inline = llm_service.describe_image_with_vision(raw_inline, mime_type=mime_inline)
                    if desc_inline:
                        image_hash_to_desc[content_hash_inline] = desc_inline
                        page_vision_blocks.setdefault(page_idx, []).append(
                            f"[PDF page {page_idx}, изображение (схема/диаграмма)]\nОписание: {desc_inline}"
                        )
                        total_described += 1
                        if progress_cb:
                            progress_cb(total_described, total_to_process)
            doc.close()
            if total_described or xref_to_desc:
                logger.info("Vision ingest (PDF): PyMuPDF found %s unique image(s), described %s", len(xref_to_desc), total_described)
                return page_vision_blocks
            logger.info("Vision ingest (PDF): PyMuPDF found no images, trying pdfplumber")
        except Exception as exc:
            logger.info("Vision ingest (PDF): PyMuPDF failed, falling back to pdfplumber: %s", exc)

    # 2) Fallback: pdfplumber (crop by bbox — may miss or crop wrong on some PDFs)
    total = 0
    total_found = 0
    with pdfplumber.open(path) as pdf:
        for page_idx, page in enumerate(pdf.pages, start=1):
            if total >= max_total:
                break
            imgs = getattr(page, "images", None) or []
            total_found += len(imgs)
            for img_idx, img_info in enumerate(imgs):
                if total >= max_total:
                    break
                try:
                    x0 = float(img_info.get("x0", 0))
                    top = float(img_info.get("top", 0))
                    x1 = float(img_info.get("x1", 0))
                    bottom = float(img_info.get("bottom", 0))
                    if x1 <= x0 or bottom <= top:
                        continue
                    cropped = page.crop((x0, top, x1, bottom))
                    pil_img = cropped.to_image(resolution=160).original
                except Exception as exc:
                    logger.debug("PDF vision: crop/render failed page %s img %s: %s", page_idx, img_idx, exc)
                    continue
                if _vision_skip_small_image(pil_img.width, pil_img.height):
                    logger.debug("PDF vision: skip small image page %s img %s (%sx%s)", page_idx, img_idx, pil_img.width, pil_img.height)
                    continue
                buf = io.BytesIO()
                try:
                    pil_img.save(buf, format="PNG")
                    raw = buf.getvalue()
                except Exception as exc:
                    logger.debug("PDF vision: save PNG failed: %s", exc)
                    continue
                desc = llm_service.describe_image_with_vision(raw, mime_type="image/png")
                if desc:
                    page_vision_blocks.setdefault(page_idx, []).append(
                        f"[PDF page {page_idx}, изображение {img_idx + 1}]\nОписание: {desc}"
                    )
                    total += 1
                    if progress_cb:
                        progress_cb(total, 0)
    if total_found or total:
        logger.info("Vision ingest (PDF): pdfplumber found %s image(s), described %s", total_found, total)
    return page_vision_blocks


def _parse_pdf(path: Path, *, progress_cb: TYPE_PROGRESS_CB | None = None) -> str:
    ocr_cfg = _ocr_runtime_settings()
    page_blocks: list[tuple[int, str]] = []
    page_table_blocks: dict[int, list[str]] = {}
    page_figure_blocks: dict[int, list[str]] = {}
    page_vision_blocks: dict[int, list[str]] = {}
    ocr_blocks: list[str] = []
    with pdfplumber.open(path) as pdf:
        for page_idx, page in enumerate(pdf.pages, start=1):
            text = page.extract_text()
            if text and text.strip():
                page_blocks.append((page_idx, text.strip()))
            table_blocks = _extract_pdf_table_blocks(page, page_idx=page_idx, page_text=text or "")
            if table_blocks:
                page_table_blocks[page_idx] = table_blocks
            if not ocr_cfg["enabled"]:
                continue
            if page_idx > int(ocr_cfg["max_pdf_pages"]):
                break
            page_images = getattr(page, "images", None) or []
            text_len = len((text or "").strip())
            needs_ocr = (not text or text_len < 160) and bool(page_images)
            if not needs_ocr:
                continue
            try:
                resolution = 220 if str(ocr_cfg["mode"]) == "accurate" else 160
                page_image = page.to_image(resolution=resolution).original
            except Exception as exc:
                logger.debug("PDF OCR render failed on page %s: %s", page_idx, exc)
                continue
            ocr_text = _ocr_image_with_tesseract(
                page_image,
                lang=str(ocr_cfg["lang"]),
                mode=str(ocr_cfg["mode"]),
            )
            cleaned = _clean_ocr_text(ocr_text, min_chars=int(ocr_cfg["min_chars"]))
            if cleaned:
                ocr_blocks.append(f"[OCR PDF page {page_idx}]\n{cleaned}")
    page_vision_blocks = _extract_pdf_vision_descriptions(path, progress_cb=progress_cb)
    cleaned_page_blocks = _prepare_pdf_page_blocks(page_blocks)
    final_page_texts: dict[int, str] = {}
    for page_idx, cleaned_text in cleaned_page_blocks:
        fig_blocks = _extract_pdf_figure_caption_blocks(cleaned_text, page_idx=page_idx)
        if fig_blocks:
            page_figure_blocks[page_idx] = fig_blocks
        captions_to_strip = _extract_captions_from_structured_blocks(page_table_blocks.get(page_idx) or [])
        captions_to_strip.extend(_extract_captions_from_structured_blocks(fig_blocks or []))
        final_page_texts[page_idx] = _strip_caption_lines_from_text_block(cleaned_text, captions_to_strip)
    # Include every page that has text, tables, figures, or vision blocks (pages with only images were missing)
    all_page_indices: set[int] = set()
    for page_idx, text in cleaned_page_blocks:
        if text and text.strip():
            all_page_indices.add(page_idx)
    for page_idx in page_vision_blocks:
        all_page_indices.add(page_idx)
    for page_idx in page_table_blocks:
        all_page_indices.add(page_idx)
    for page_idx in page_figure_blocks:
        all_page_indices.add(page_idx)
    pages = [
        (
            f"[PDF page {page_idx}]\n{(final_page_texts.get(page_idx) or '')}"
            + (
                "\n\n" + "\n\n".join(page_table_blocks.get(page_idx) or [])
                if page_table_blocks.get(page_idx)
                else ""
            )
            + (
                "\n\n" + "\n\n".join(page_figure_blocks.get(page_idx) or [])
                if page_figure_blocks.get(page_idx)
                else ""
            )
            + (
                "\n\n" + "\n\n".join(page_vision_blocks.get(page_idx) or [])
                if page_vision_blocks.get(page_idx)
                else ""
            )
        )
        for page_idx in sorted(all_page_indices)
    ]
    if ocr_blocks:
        pages.append("\n\n".join(ocr_blocks))
    return "\n\n".join(pages)


def _parse_docx(path: Path, *, progress_cb: TYPE_PROGRESS_CB | None = None) -> str:
    ocr_cfg = _ocr_runtime_settings()
    vision_cfg = get_vision_ingest_settings()
    doc = DocxDocument(str(path))
    parts: list[str] = []
    for p in doc.paragraphs:
        txt = str(getattr(p, "text", "") or "").strip()
        if not txt:
            continue
        style_name = str(getattr(getattr(p, "style", None), "name", "") or "")
        heading_level = _docx_heading_level(style_name)
        if heading_level:
            parts.append(f"{'#' * heading_level} {txt}")
        else:
            parts.append(txt)
    if ocr_cfg["enabled"] or vision_cfg.get("enabled"):
        ocr_blocks = _extract_docx_image_ocr(path, ocr_cfg=ocr_cfg, progress_cb=progress_cb)
        if ocr_blocks:
            parts.append("\n\n".join(ocr_blocks))
    return _inject_section_anchor_markers("\n\n".join(parts), scope_prefix="docx")


def _extract_pptx_vision_descriptions(
    prs: Presentation,
    *,
    progress_cb: TYPE_PROGRESS_CB | None = None,
) -> dict[int, list[str]]:
    """Extract picture shapes from PPTX slides, describe via vision LLM; return slide_idx -> list of blocks."""
    vision_cfg = get_vision_ingest_settings()
    if not vision_cfg.get("enabled"):
        return {}
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except (ImportError, AttributeError):
        return {}
    max_total = max(1, int(vision_cfg.get("max_images_per_document", 20)))
    slide_vision_blocks: dict[int, list[str]] = {}
    total = 0
    try:
        from app.services import llm_service
    except ImportError:
        return {}
    for slide_idx, slide in enumerate(prs.slides, start=1):
        if total >= max_total:
            break
        for shape in getattr(slide, "shapes", []) or []:
            if total >= max_total:
                break
            try:
                if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE:
                    continue
                img = getattr(shape, "image", None)
                if not img:
                    continue
                raw = getattr(img, "blob", None)
                if not raw:
                    continue
                content_type = getattr(img, "content_type", None) or "image/png"
                if content_type and "/" in str(content_type):
                    mime = str(content_type).strip()
                else:
                    mime = _mime_for_image_path(f".{getattr(img, 'ext', 'png')}")
            except Exception as exc:
                logger.debug("PPTX vision: get image failed slide %s: %s", slide_idx, exc)
                continue
            try:
                pil_check = Image.open(io.BytesIO(raw))
                pw, ph = pil_check.size
                pil_check.close()
            except (UnidentifiedImageError, OSError):
                pw, ph = 0, 0
            if _vision_skip_small_image(pw, ph):
                logger.debug("PPTX vision: skip small image slide %s (%sx%s)", slide_idx, pw, ph)
                continue
            desc = llm_service.describe_image_with_vision(raw, mime_type=mime)
            if desc:
                slide_vision_blocks.setdefault(slide_idx, []).append(
                    f"[PPTX слайд {slide_idx}, изображение]\nОписание: {desc}"
                )
                total += 1
                if progress_cb:
                    progress_cb(total, 0)
    return slide_vision_blocks


def _pptx_fallback_text_from_zip(path: Path) -> str:
    """Extract slide text from PPTX as zip when python-pptx fails (e.g. SVG without content-type)."""
    parts: list[str] = []
    try:
        with zipfile.ZipFile(path, "r") as zf:
            slide_names = sorted(n for n in zf.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml"))
            for idx, name in enumerate(slide_names, start=1):
                try:
                    xml = zf.read(name).decode("utf-8", errors="replace")
                except Exception:
                    continue
                # OOXML text in <a:t>...</a:t>
                texts = re.findall(r"<a:t>([^<]*)</a:t>", xml)
                line = " ".join(t.strip() for t in texts if t.strip())
                if line:
                    parts.append(f"Слайд {idx}\n{line}")
    except Exception as exc:
        logger.debug("PPTX zip fallback failed: %s", exc)
    return "\n\n".join(parts) if parts else ""


def _parse_pptx(path: Path, *, progress_cb: TYPE_PROGRESS_CB | None = None) -> str:
    try:
        prs = Presentation(str(path))
    except KeyError as e:
        logger.warning("PPTX open failed (e.g. SVG without content-type): %s. Using zip fallback.", e)
        fallback = _pptx_fallback_text_from_zip(path)
        if fallback:
            return fallback
        raise ValueError(
            "Презентация содержит неподдерживаемые элементы (например SVG без типа контента). "
            "Пересохраните файл в PowerPoint или экспортируйте в PDF."
        ) from e
    except Exception as e:
        logger.warning("PPTX open failed: %s. Using zip fallback.", e)
        fallback = _pptx_fallback_text_from_zip(path)
        if fallback:
            return fallback
        raise ValueError(f"Не удалось открыть презентацию: {e!s}") from e

    slide_vision_blocks = _extract_pptx_vision_descriptions(prs, progress_cb=progress_cb)
    slides_text: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        table_blocks: list[str] = []
        _collect_pptx_shape_blocks(
            getattr(slide, "shapes", []) or [],
            slide_idx=idx,
            text_parts=parts,
            table_blocks=table_blocks,
            table_counter=[0],
        )
        figure_blocks = _extract_pptx_figure_caption_blocks(parts, slide_idx=idx)
        captions_to_strip = _extract_captions_from_structured_blocks(table_blocks)
        captions_to_strip.extend(_extract_captions_from_structured_blocks(figure_blocks))
        filtered_parts = _filter_text_parts_remove_captions(parts, captions_to_strip)
        plain_text_block = _inject_section_anchor_markers(
            _annotate_heading_markers("\n".join(filtered_parts)),
            scope_prefix=f"pptx:s{int(idx)}",
        )
        combined_parts: list[str] = []
        if plain_text_block:
            combined_parts.append(plain_text_block)
        if table_blocks:
            combined_parts.extend(table_blocks)
        if figure_blocks:
            combined_parts.extend(figure_blocks)
        for block in slide_vision_blocks.get(idx) or []:
            combined_parts.append(block)
        if combined_parts:
            slides_text.append(f"Слайд {idx}\n" + "\n\n".join(combined_parts))
    return "\n\n".join(slides_text)


def _parse_html(html: str) -> str:
    text = trafilatura.extract(html)
    return text or ""


def _tesseract_available() -> bool:
    global _TESSERACT_OK
    if _TESSERACT_OK is not None:
        return _TESSERACT_OK
    _TESSERACT_OK = shutil.which("tesseract") is not None
    if not _TESSERACT_OK:
        logger.info("Tesseract is not installed; OCR for document images is disabled.")
    return _TESSERACT_OK


def _clean_ocr_text(raw: str | None, *, min_chars: int = OCR_DEFAULT_MIN_CHARS) -> str:
    txt = (raw or "").strip()
    if not txt:
        return ""
    txt = "\n".join(line.strip() for line in txt.splitlines())
    txt = "\n".join(line for line in txt.splitlines() if line)
    if len(txt) < max(1, int(min_chars)):
        return ""
    return txt


def _docx_heading_level(style_name: str | None) -> int | None:
    style = str(style_name or "").strip()
    if not style:
        return None
    m = _DOCX_HEADING_STYLE_RE.search(style)
    if not m:
        return None
    try:
        lvl = int(m.group(1) or "1")
    except (TypeError, ValueError):
        lvl = 1
    return max(1, min(6, lvl))


def _normalize_ws(line: str) -> str:
    return re.sub(r"\s+", " ", str(line or "")).strip()


def _normalize_table_cell(value: object) -> str:
    txt = str(value or "").replace("\n", " ").replace("\r", " ")
    return _normalize_ws(txt)


def _match_key_loose(text: str) -> str:
    txt = _normalize_ws(text).lower()
    txt = re.sub(r"[^a-zа-яё0-9]+", "", txt, flags=re.IGNORECASE)
    return txt


def _inject_section_anchor_markers(text: str | None, *, scope_prefix: str) -> str:
    src = str(text or "").strip()
    prefix = str(scope_prefix or "").strip()
    if not src or not prefix:
        return src
    lines = src.splitlines()
    out: list[str] = []
    sec_idx = 0
    for raw in lines:
        line = raw.rstrip()
        stripped = line.strip()
        md = _MD_HEADING_LINE_RE.match(stripped)
        if md:
            prev_nonempty = ""
            for prev in reversed(out):
                if prev.strip():
                    prev_nonempty = prev.strip()
                    break
            if not _ANCHOR_LINE_RE.match(prev_nonempty):
                sec_idx += 1
                out.append(f"Якорь: {prefix}:sec:{sec_idx}")
        out.append(line)
    return "\n".join(out).strip()


def _is_page_number_line(line: str) -> bool:
    txt = _normalize_ws(line).lower()
    if not txt:
        return False
    return bool(
        re.fullmatch(r"(?:стр(?:аница)?\.?\s*)?\d{1,4}(?:\s*(?:/|из|of)\s*\d{1,4})?", txt)
        or re.fullmatch(r"page\s+\d{1,4}(?:\s*(?:/|of)\s*\d{1,4})?", txt)
    )


def _pdf_margin_line_key(line: str) -> str:
    txt = _normalize_ws(line).lower()
    if not txt:
        return ""
    txt = re.sub(r"\d+", "#", txt)
    txt = re.sub(r"[^\w\s#./-]+", "", txt, flags=re.UNICODE)
    txt = re.sub(r"\s+", " ", txt).strip()
    return txt


def _is_pdf_margin_candidate(line: str) -> bool:
    txt = _normalize_ws(line)
    if not txt:
        return False
    if _PDF_MARKER_RE.match(txt):
        return False
    if len(txt) > PDF_MARGIN_CANDIDATE_MAX_LEN:
        return False
    if txt.startswith(("-", "*", "•")):
        return False
    word_count = len(txt.split())
    if word_count > 16:
        return False
    if not re.search(r"[A-Za-zА-Яа-яЁё0-9]", txt):
        return False
    return True


def _strip_repeated_pdf_margin_noise(page_texts: list[str]) -> tuple[list[str], int]:
    if len(page_texts) < PDF_MARGIN_REPEAT_MIN_PAGES:
        return page_texts, 0
    line_pages: dict[str, set[int]] = defaultdict(set)
    page_lines: list[list[str]] = []
    for page_idx, page_text in enumerate(page_texts):
        lines = [ln.rstrip() for ln in str(page_text or "").splitlines()]
        page_lines.append(lines)
        if not lines:
            continue
        margin_slice = lines[:PDF_MARGIN_SCAN_LINES] + lines[-PDF_MARGIN_SCAN_LINES:]
        for line in margin_slice:
            if not _is_pdf_margin_candidate(line):
                continue
            key = _pdf_margin_line_key(line)
            if not key:
                continue
            line_pages[key].add(page_idx)
    repeated_keys = {
        key
        for key, pages in line_pages.items()
        if len(pages) >= PDF_MARGIN_REPEAT_MIN_PAGES
    }
    if not repeated_keys:
        return page_texts, 0

    removed = 0
    cleaned_pages: list[str] = []
    for lines in page_lines:
        if not lines:
            cleaned_pages.append("")
            continue
        last_idx = len(lines) - 1
        next_lines: list[str] = []
        for i, raw_line in enumerate(lines):
            line = raw_line.rstrip()
            is_margin_zone = i < PDF_MARGIN_SCAN_LINES or i > last_idx - PDF_MARGIN_SCAN_LINES
            if not is_margin_zone:
                next_lines.append(line)
                continue
            key = _pdf_margin_line_key(line)
            drop = False
            if key and key in repeated_keys and _is_pdf_margin_candidate(line):
                drop = True
            elif _is_page_number_line(line) and len(lines) > 2:
                drop = True
            if drop:
                removed += 1
                continue
            next_lines.append(line)
        cleaned_pages.append("\n".join(next_lines).strip())
    return cleaned_pages, removed


def _looks_like_heading_line(line: str, *, prev_blank: bool, next_blank: bool) -> bool:
    txt = _normalize_ws(line)
    if not txt:
        return False
    if txt.startswith("[") or _PDF_MARKER_RE.match(txt):
        return False
    if txt.startswith(("-", "*", "•")):
        return False
    if txt.startswith("#"):
        return False
    if _is_page_number_line(txt):
        return False
    if not re.search(r"[A-Za-zА-Яа-яЁё]", txt):
        return False
    if len(txt) < 4 or len(txt) > 130:
        return False
    if txt.endswith((".", "!", "?")):
        return False
    if txt.count(",") > 1:
        return False
    words = txt.split()
    if not (1 <= len(words) <= 12):
        return False

    # Strong heading signals that don't require blank separators.
    letters = [ch for ch in txt if ch.isalpha()]
    uppercase_ratio = (sum(1 for ch in letters if ch.isupper()) / len(letters)) if letters else 0.0
    strong_signal = uppercase_ratio >= 0.65 or bool(re.match(r"^[IVXLCM]+\.\s+\S+", txt, flags=re.IGNORECASE))
    if not (prev_blank or next_blank or strong_signal):
        return False

    # Avoid tagging likely wrapped sentence fragments.
    if txt[0].islower():
        return False
    if len(words) >= 6 and uppercase_ratio < 0.22 and ":" not in txt:
        return False
    return True


def _annotate_heading_markers(text: str) -> str:
    if not text:
        return ""
    lines = str(text).splitlines()
    if not lines:
        return ""
    out: list[str] = []
    for i, raw in enumerate(lines):
        line = raw.rstrip()
        stripped = line.strip()
        prev_blank = i == 0 or not lines[i - 1].strip()
        next_blank = i == len(lines) - 1 or not lines[i + 1].strip()
        if _looks_like_heading_line(stripped, prev_blank=prev_blank, next_blank=next_blank):
            hashes = "# " if sum(1 for ch in stripped if ch.isalpha() and ch.isupper()) >= max(3, len([c for c in stripped if c.isalpha()]) * 0.65) else "## "
            out.append(f"{hashes}{stripped}")
        else:
            out.append(line)
    return "\n".join(out).strip()


def _prepare_pdf_page_blocks(page_blocks: list[tuple[int, str]]) -> list[tuple[int, str]]:
    if not page_blocks:
        return []
    page_ids = [page_idx for page_idx, _ in page_blocks]
    raw_texts = [_normalize_ws_lines_preserve_breaks(text) for _, text in page_blocks]
    cleaned_texts, removed = _strip_repeated_pdf_margin_noise(raw_texts)
    if removed:
        logger.info("PDF ingest cleanup: removed %d repeated margin/page-noise lines", removed)
    out: list[tuple[int, str]] = []
    for page_idx, txt in zip(page_ids, cleaned_texts):
        normalized = _normalize_ws_lines_preserve_breaks(txt)
        if not normalized:
            continue
        with_headings = _annotate_heading_markers(normalized)
        with_anchors = _inject_section_anchor_markers(with_headings, scope_prefix=f"pdf:p{int(page_idx)}")
        out.append((page_idx, with_anchors))
    return out


def _normalize_ws_lines_preserve_breaks(text: str | None) -> str:
    lines = [_normalize_ws(ln) for ln in str(text or "").splitlines()]
    # Keep intentional blank separators, but collapse duplicate blank runs.
    out: list[str] = []
    prev_blank = False
    for line in lines:
        is_blank = not line
        if is_blank and prev_blank:
            continue
        out.append(line)
        prev_blank = is_blank
    return "\n".join(out).strip()


def _dedupe_text_lines_preserve_order(lines: list[str]) -> list[str]:
    seen: set[str] = set()
    out: list[str] = []
    for raw in lines or []:
        line = str(raw or "").strip()
        if not line:
            continue
        key = _match_key_loose(line)
        if not key:
            continue
        if key in seen:
            continue
        seen.add(key)
        out.append(line)
    return out


def _find_caption_candidates(text: str | None) -> dict[str, list[str]]:
    out = {"table": [], "figure": []}
    for raw_line in str(text or "").splitlines():
        line = _normalize_ws(raw_line)
        if not line or len(line) < 6 or len(line) > 220:
            continue
        if _TABLE_CAPTION_RE.match(line):
            out["table"].append(line)
        elif _FIGURE_CAPTION_RE.match(line):
            out["figure"].append(line)
    out["table"] = _dedupe_text_lines_preserve_order(out["table"])
    out["figure"] = _dedupe_text_lines_preserve_order(out["figure"])
    return out


def _caption_mentions_index(caption: str, index: int) -> bool:
    if not caption or index <= 0:
        return False
    nums = re.findall(r"\d+(?:\.\d+)?", caption)
    if any(n.split(".", 1)[0] == str(index) for n in nums):
        return True
    roman = {1: "i", 2: "ii", 3: "iii", 4: "iv", 5: "v", 6: "vi", 7: "vii", 8: "viii", 9: "ix", 10: "x"}
    roman_token = roman.get(index)
    if roman_token and re.search(rf"\b{re.escape(roman_token)}\b", caption, flags=re.IGNORECASE):
        return True
    return False


def _pick_table_caption(page_text: str, table_idx: int, used: set[str] | None = None) -> str:
    candidates = _find_caption_candidates(page_text).get("table") or []
    if not candidates:
        return ""
    used_keys = used if isinstance(used, set) else set()
    for caption in candidates:
        key = _match_key_loose(caption)
        if key in used_keys:
            continue
        if _caption_mentions_index(caption, int(table_idx)):
            used_keys.add(key)
            return caption
    for caption in candidates:
        key = _match_key_loose(caption)
        if key in used_keys:
            continue
        used_keys.add(key)
        return caption
    return ""


def _extract_figure_caption_blocks_from_text(
    text: str | None,
    *,
    source: str,
    page_idx: int | None = None,
    slide_idx: int | None = None,
    max_items: int = 8,
) -> list[str]:
    captions = (_find_caption_candidates(text).get("figure") or [])[: max(0, int(max_items or 0))]
    blocks: list[str] = []
    seen: set[str] = set()
    for fig_idx, caption in enumerate(captions, start=1):
        cap = _normalize_ws(caption)
        if not cap:
            continue
        key = _match_key_loose(cap)
        if key and key in seen:
            continue
        if key:
            seen.add(key)
        if source == "pdf" and page_idx is not None:
            lines = [
                f"[PDF page {int(page_idx)}]",
                f"[PDF figure {int(fig_idx)}]",
                f"Якорь: pdf:p{int(page_idx)}:fig:{int(fig_idx)}",
                f"## Рисунок {int(fig_idx)} (страница {int(page_idx)})",
                f"Подпись: {cap}",
            ]
        elif source == "pptx" and slide_idx is not None:
            lines = [
                f"Слайд {int(slide_idx)}",
                f"[PPTX figure {int(fig_idx)}]",
                f"Якорь: pptx:s{int(slide_idx)}:fig:{int(fig_idx)}",
                f"## Рисунок {int(fig_idx)} (слайд {int(slide_idx)})",
                f"Подпись: {cap}",
            ]
        else:
            continue
        blocks.append("\n".join(lines).strip())
    return blocks


def _extract_pdf_figure_caption_blocks(page_text: str, *, page_idx: int) -> list[str]:
    return _extract_figure_caption_blocks_from_text(
        page_text,
        source="pdf",
        page_idx=page_idx,
        max_items=PDF_FIGURE_CAPTION_MAX_PER_PAGE,
    )


def _extract_pptx_figure_caption_blocks(text_parts: list[str], *, slide_idx: int) -> list[str]:
    joined = "\n".join(str(x or "") for x in (text_parts or []))
    return _extract_figure_caption_blocks_from_text(
        joined,
        source="pptx",
        slide_idx=slide_idx,
        max_items=PPTX_FIGURE_CAPTION_MAX_PER_SLIDE,
    )


def _extract_captions_from_structured_blocks(blocks: list[str]) -> list[str]:
    out: list[str] = []
    seen: set[str] = set()
    for block in blocks or []:
        for raw_line in str(block or "").splitlines():
            m = _CAPTION_LINE_RE.match(raw_line.strip())
            if not m:
                continue
            caption = _normalize_ws(m.group(1))
            if not caption:
                continue
            key = _match_key_loose(caption)
            if key and key in seen:
                continue
            if key:
                seen.add(key)
            out.append(caption)
    return out


def _strip_caption_lines_from_text_block(text: str | None, captions: list[str]) -> str:
    src = str(text or "")
    if not src.strip():
        return ""
    caption_keys = {_match_key_loose(c) for c in (captions or []) if _match_key_loose(c)}
    if not caption_keys:
        return src.strip()
    kept: list[str] = []
    for raw_line in src.splitlines():
        line = raw_line.rstrip()
        key = _match_key_loose(line)
        if key and key in caption_keys:
            continue
        kept.append(line)
    return _normalize_ws_lines_preserve_breaks("\n".join(kept))


def _filter_text_parts_remove_captions(parts: list[str], captions: list[str]) -> list[str]:
    caption_keys = {_match_key_loose(c) for c in (captions or []) if _match_key_loose(c)}
    if not caption_keys:
        return list(parts or [])
    out: list[str] = []
    for part in parts or []:
        text = str(part or "").strip()
        if not text:
            continue
        key = _match_key_loose(text)
        if key and key in caption_keys:
            continue
        out.append(text)
    return out


def _table_block_dedupe_key(serialized: str) -> str:
    if not serialized:
        return ""
    useful_lines = []
    for line in str(serialized).splitlines():
        if line.startswith("[PDF page ") or line.startswith("[PDF table ") or line.startswith("[PPTX table "):
            continue
        if line.startswith("Слайд ") or line.startswith("Якорь: ") or line.startswith("Подпись: "):
            continue
        if line.startswith("## Таблица "):
            continue
        useful_lines.append(line)
    return _match_key_loose("\n".join(useful_lines))


def _trim_empty_table_columns(rows: list[list[str]]) -> list[list[str]]:
    if not rows:
        return rows
    width = max((len(r) for r in rows), default=0)
    if width <= 0:
        return []
    padded = [r + [""] * (width - len(r)) for r in rows]
    keep_cols = [
        idx for idx in range(width)
        if any(_normalize_ws(r[idx]) for r in padded)
    ]
    if not keep_cols:
        return []
    return [[_normalize_ws(r[idx]) for idx in keep_cols] for r in padded]


def _table_block_redundant(serialized: str, page_text: str) -> bool:
    page_norm = _normalize_ws(str(page_text or "")).lower()
    page_key = _match_key_loose(page_norm)
    if not serialized or not page_norm:
        return False
    row_lines = [
        _normalize_ws(ln.split(":", 1)[1] if ":" in ln else ln).lower()
        for ln in serialized.splitlines()
        if ln.startswith("Строка ")
    ]
    row_lines = [ln for ln in row_lines if ln]
    if not row_lines:
        return False
    hits = 0
    for ln in row_lines:
        if ln in page_norm:
            hits += 1
            continue
        row_key = _match_key_loose(ln)
        if row_key and row_key in page_key:
            hits += 1
    return hits >= max(2, int(len(row_lines) * 0.7))


def _serialize_pdf_table(
    table: list[list[object]],
    *,
    page_idx: int,
    table_idx: int,
    page_text: str = "",
    caption: str = "",
) -> str:
    if not isinstance(table, list) or not table:
        return ""
    norm_rows: list[list[str]] = []
    for row in table[:PDF_TABLE_MAX_ROWS]:
        if not isinstance(row, list):
            continue
        cells = [_normalize_table_cell(cell) for cell in row[:PDF_TABLE_MAX_COLS]]
        if any(cells):
            norm_rows.append(cells)
    if len(norm_rows) < 2:
        return ""
    norm_rows = _trim_empty_table_columns(norm_rows)
    if not norm_rows:
        return ""
    filled_cells = sum(1 for row in norm_rows for c in row if c)
    if filled_cells < PDF_TABLE_MIN_FILLED_CELLS:
        return ""
    col_count = max((len(r) for r in norm_rows), default=0)
    if col_count < 2:
        return ""

    def _cell_join(cells: list[str]) -> str:
        return " | ".join((c if c else "—") for c in cells)

    header = norm_rows[0]
    data_rows = norm_rows[1:] if len(norm_rows) > 1 else []
    lines = [
        f"[PDF page {int(page_idx)}]",
        f"[PDF table {int(table_idx)}]",
        f"Якорь: pdf:p{int(page_idx)}:table:{int(table_idx)}",
        f"## Таблица {int(table_idx)} (страница {int(page_idx)})",
    ]
    caption_clean = _normalize_ws(caption)
    if caption_clean:
        lines.append(f"Подпись: {caption_clean}")
    lines.append(f"Колонки: {_cell_join(header)}")
    for idx, row in enumerate(data_rows, start=1):
        lines.append(f"Строка {idx}: {_cell_join(row)}")
    serialized = "\n".join(lines).strip()
    if _table_block_redundant(serialized, page_text):
        return ""
    return serialized


def _extract_pdf_table_blocks(page, *, page_idx: int, page_text: str) -> list[str]:
    try:
        raw_tables = page.extract_tables() or []
    except Exception as exc:
        logger.debug("PDF table extraction failed on page %s: %s", page_idx, exc)
        return []
    blocks: list[str] = []
    seen_keys: set[str] = set()
    used_captions: set[str] = set()
    for t_idx, table in enumerate(raw_tables[:PDF_TABLE_MAX_PER_PAGE], start=1):
        caption = _pick_table_caption(page_text, t_idx, used=used_captions)
        block = _serialize_pdf_table(
            table,
            page_idx=page_idx,
            table_idx=t_idx,
            page_text=page_text,
            caption=caption,
        )
        if block:
            key = _table_block_dedupe_key(block)
            if key and key in seen_keys:
                continue
            if key:
                seen_keys.add(key)
            blocks.append(block)
    return blocks


def _serialize_pptx_table(
    table_rows: list[list[object]],
    *,
    slide_idx: int,
    table_idx: int,
    caption: str = "",
) -> str:
    norm_rows: list[list[str]] = []
    for row in (table_rows or [])[:PPTX_TABLE_MAX_ROWS]:
        if not isinstance(row, list):
            continue
        cells = [_normalize_table_cell(cell) for cell in row[:PPTX_TABLE_MAX_COLS]]
        if any(cells):
            norm_rows.append(cells)
    if len(norm_rows) < 2:
        return ""
    norm_rows = _trim_empty_table_columns(norm_rows)
    if not norm_rows:
        return ""
    if sum(1 for r in norm_rows for c in r if c) < PDF_TABLE_MIN_FILLED_CELLS:
        return ""
    if max((len(r) for r in norm_rows), default=0) < 2:
        return ""

    def _join_cells(cells: list[str]) -> str:
        return " | ".join((c if c else "—") for c in cells)

    header = norm_rows[0]
    data_rows = norm_rows[1:]
    lines = [
        f"Слайд {int(slide_idx)}",
        f"[PPTX table {int(table_idx)}]",
        f"Якорь: pptx:s{int(slide_idx)}:table:{int(table_idx)}",
        f"## Таблица {int(table_idx)} (слайд {int(slide_idx)})",
    ]
    caption_clean = _normalize_ws(caption)
    if caption_clean:
        lines.append(f"Подпись: {caption_clean}")
    lines.append(f"Колонки: {_join_cells(header)}")
    for i, row in enumerate(data_rows, start=1):
        lines.append(f"Строка {i}: {_join_cells(row)}")
    return "\n".join(lines).strip()


def _pptx_shape_table_rows(shape) -> list[list[object]]:
    if not bool(getattr(shape, "has_table", False)):
        return []
    table = getattr(shape, "table", None)
    if table is None:
        return []
    out: list[list[object]] = []
    try:
        rows = getattr(table, "rows", None) or []
        for row in rows:
            cells = getattr(row, "cells", None) or []
            out.append([getattr(cell, "text", "") for cell in cells])
    except Exception:
        return []
    return out


def _collect_pptx_shape_blocks(
    shapes,
    *,
    slide_idx: int,
    text_parts: list[str],
    table_blocks: list[str],
    table_counter: list[int],
    recent_texts: list[str] | None = None,
    seen_text_keys: set[str] | None = None,
    seen_table_keys: set[str] | None = None,
) -> None:
    recent = recent_texts if isinstance(recent_texts, list) else []
    seen_text = seen_text_keys if isinstance(seen_text_keys, set) else set()
    seen_tables = seen_table_keys if isinstance(seen_table_keys, set) else set()
    for shape in shapes or []:
        # Group shapes may contain nested shapes.
        child_shapes = getattr(shape, "shapes", None)
        has_table = bool(getattr(shape, "has_table", False))
        if has_table:
            table_rows = _pptx_shape_table_rows(shape)
            if table_rows:
                if table_counter:
                    table_counter[0] = int(table_counter[0]) + 1
                    t_idx = int(table_counter[0])
                else:
                    t_idx = 1
                if t_idx <= PPTX_TABLE_MAX_PER_SLIDE:
                    caption = ""
                    for candidate in reversed(recent[-4:]):
                        if _TABLE_CAPTION_RE.match(candidate):
                            caption = candidate
                            break
                    block = _serialize_pptx_table(table_rows, slide_idx=slide_idx, table_idx=t_idx, caption=caption)
                    if block:
                        table_key = _table_block_dedupe_key(block)
                        if not table_key or table_key not in seen_tables:
                            if table_key:
                                seen_tables.add(table_key)
                            table_blocks.append(block)
            # Do not also take shape.text for table shapes: it usually duplicates cells.
            continue
        text = getattr(shape, "text", None)
        if isinstance(text, str):
            cleaned = text.strip()
            if cleaned:
                key = _match_key_loose(cleaned)
                if key and key not in seen_text:
                    seen_text.add(key)
                    text_parts.append(cleaned)
                recent.append(cleaned)
                if len(recent) > 8:
                    del recent[:-8]
        if child_shapes:
            _collect_pptx_shape_blocks(
                child_shapes,
                slide_idx=slide_idx,
                text_parts=text_parts,
                table_blocks=table_blocks,
                table_counter=table_counter,
                recent_texts=recent,
                seen_text_keys=seen_text,
                seen_table_keys=seen_tables,
            )


def _ocr_image_with_tesseract(
    image: Image.Image,
    *,
    lang: str = OCR_DEFAULT_LANG,
    mode: str = OCR_DEFAULT_MODE,
) -> str:
    if not _tesseract_available():
        return ""
    try:
        rgb = image.convert("RGB")
    except Exception:
        return ""
    with tempfile.TemporaryDirectory(prefix="ocr_") as tmp:
        inp = Path(tmp) / "img.png"
        out_base = Path(tmp) / "out"
        try:
            rgb.save(inp, format="PNG")
        except Exception:
            return ""
        try:
            psm = "3" if str(mode).strip().lower() == "accurate" else "6"
            proc = subprocess.run(
                ["tesseract", str(inp), str(out_base), "-l", str(lang or OCR_DEFAULT_LANG), "--oem", "1", "--psm", psm],
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                check=False,
                timeout=30,
            )
        except Exception as exc:
            logger.debug("OCR execution failed: %s", exc)
            return ""
        if proc.returncode != 0:
            logger.debug("Tesseract returned %s: %s", proc.returncode, (proc.stderr or "").strip())
            return ""
        out_txt = out_base.with_suffix(".txt")
        if not out_txt.exists():
            return ""
        try:
            return out_txt.read_text(encoding="utf-8", errors="replace")
        except Exception:
            return ""


def _mime_for_image_path(name: str) -> str:
    ext = (Path(name).suffix or "").lower()
    if ext in (".png",):
        return "image/png"
    if ext in (".jpg", ".jpeg",):
        return "image/jpeg"
    if ext in (".webp",):
        return "image/webp"
    return "image/png"


def _docx_collect_image_parts(part: object, seen: set, out: list[tuple[bytes, str]]) -> None:
    """Recursively collect image blobs from document part and related_parts (python-docx)."""
    try:
        related = getattr(part, "related_parts", None) or {}
    except Exception:
        return
    for _rId, rel_part in related.items():
        try:
            part_id = id(rel_part)
            if part_id in seen:
                continue
            seen.add(part_id)
            ct = (getattr(rel_part, "content_type", None) or "").strip().lower()
            if ct.startswith("image/"):
                blob = getattr(rel_part, "blob", None) or getattr(rel_part, "_blob", None)
                if blob and len(blob) > 0:
                    out.append((bytes(blob), ct))
            _docx_collect_image_parts(rel_part, seen, out)
        except Exception:
            continue


def _extract_docx_image_ocr(
    path: Path,
    *,
    ocr_cfg: dict | None = None,
    progress_cb: TYPE_PROGRESS_CB | None = None,
) -> list[str]:
    cfg = ocr_cfg or _ocr_runtime_settings()
    max_images = max(1, int(cfg.get("max_docx_images", OCR_DEFAULT_MAX_DOCX_IMAGES)))
    lang = str(cfg.get("lang", OCR_DEFAULT_LANG))
    mode = str(cfg.get("mode", OCR_DEFAULT_MODE))
    min_chars = max(1, int(cfg.get("min_chars", OCR_DEFAULT_MIN_CHARS)))
    vision_cfg = get_vision_ingest_settings()
    vision_enabled = bool(vision_cfg.get("enabled"))
    vision_max = max(1, min(max_images, int(vision_cfg.get("max_images_per_document", 20)))) if vision_enabled else 0
    if vision_enabled:
        logger.info("Vision ingest (DOCX): enabled, max %s images", vision_max)

    blocks: list[str] = []
    # Collect (raw_bytes, mime) from document part relations first, then zip
    image_items: list[tuple[bytes, str]] = []
    try:
        doc = DocxDocument(str(path))
        _docx_collect_image_parts(doc.part, set(), image_items)
    except Exception as exc:
        logger.debug("DOCX part images: %s", exc)
    if not image_items:
        image_exts = (".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tiff", ".tif")
        try:
            with zipfile.ZipFile(path, "r") as zf:
                media_files = [n for n in zf.namelist() if n.startswith("word/media/")]
                if not media_files:
                    media_files = [
                        n for n in zf.namelist()
                        if Path(n).suffix.lower() in image_exts and not n.startswith("__")
                    ]
                for name in media_files[:max_images]:
                    try:
                        raw = zf.read(name)
                        mime = _mime_for_image_path(name)
                        image_items.append((raw, mime))
                    except Exception:
                        continue
        except Exception as exc:
            logger.debug("DOCX zip images: %s", exc)
    if vision_enabled and image_items:
        logger.info("Vision ingest (DOCX): found %s image(s) in %s", len(image_items), path.name)
    try:
        for idx, (raw, mime) in enumerate(image_items[:max_images], start=1):
            try:
                img = Image.open(io.BytesIO(raw))
            except (UnidentifiedImageError, OSError):
                continue
            img_w, img_h = img.size
            ocr_text = ""
            if cfg.get("enabled", True):
                try:
                    ocr_text = _clean_ocr_text(
                        _ocr_image_with_tesseract(img, lang=lang, mode=mode),
                        min_chars=min_chars,
                    )
                finally:
                    try:
                        img.close()
                    except Exception:
                        pass
            else:
                try:
                    img.close()
                except Exception:
                    pass

            vision_desc = None
            if vision_enabled and idx <= vision_max and not _vision_skip_small_image(img_w, img_h):
                try:
                    from app.services import llm_service
                    vision_desc = llm_service.describe_image_with_vision(raw, mime_type=mime)
                    if vision_desc is None and ocr_text:
                        logger.info("Vision ingest (DOCX): image %s — LLM не вернул описание, в блоке только OCR", idx)
                except Exception as exc:
                    logger.debug("Vision description for DOCX image %s failed: %s", idx, exc)

            if not ocr_text and not vision_desc:
                continue
            # Сначала описание от vision (основное), затем OCR с подписью, чтобы оба не терялись
            parts = [f"[DOCX image {idx}]"]
            if vision_desc:
                logger.info("Vision ingest: adding description to DOCX image %s (%s chars)", idx, len(vision_desc))
                parts.append(f"Описание изображения: {vision_desc}")
            if ocr_text:
                parts.append(f"Текст с изображения (OCR): {ocr_text}")
            blocks.append("\n".join(parts))
            if progress_cb:
                progress_cb(idx, len(image_items))
    except Exception as exc:
        logger.debug("DOCX image OCR failed: %s", exc)
    return blocks
